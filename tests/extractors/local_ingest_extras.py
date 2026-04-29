"""local_ingest_extras — cheap-stage probe for non-PDF formats.

Measures whether stdlib `csv` + `openpyxl` + `python-pptx` would close
the gaps that pypdf-only ingest leaves on csv/xlsx/pptx fixtures. Emits
GFM markdown for each.

These extractors are NOT wired into local_ingest.py this round — they
live in the harness purely as a comparison probe.
"""
from pathlib import Path
import io

NAME = "local_ingest_extras"
SUPPORTS = {"csv", "xlsx", "pptx"}


def available() -> tuple[bool, str]:
    missing = []
    try:
        import openpyxl  # noqa: F401
    except Exception:
        missing.append("openpyxl")
    try:
        import pptx  # noqa: F401
    except Exception:
        missing.append("python-pptx")
    if missing:
        return False, f"missing: {', '.join(missing)}"
    return True, ""


def _gfm(rows: list) -> str:
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    norm = [
        [("" if c is None else str(c).replace("|", "\\|").replace("\n", " "))
         for c in r + [""] * (width - len(r))]
        for r in rows
    ]
    out = ["| " + " | ".join(norm[0]) + " |",
           "|" + "|".join(["---"] * width) + "|"]
    for r in norm[1:]:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)


def _extract_csv(raw: bytes) -> str:
    import csv
    text = raw.decode("utf-8", errors="replace")
    rows = list(csv.reader(io.StringIO(text)))
    return _gfm(rows)


def _extract_xlsx(raw: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(raw), data_only=True, read_only=True)
    chunks = []
    for sheet in wb.worksheets:
        rows = [list(r) for r in sheet.iter_rows(values_only=True)]
        if not rows:
            continue
        chunks.append(f"## Sheet: {sheet.title}\n")
        chunks.append(_gfm(rows))
    return "\n".join(chunks)


def _extract_pptx(raw: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(raw))
    chunks = []
    for i, slide in enumerate(prs.slides, 1):
        chunks.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    chunks.append(t)
            if shape.has_table:
                rows = [[cell.text.strip() for cell in row.cells]
                        for row in shape.table.rows]
                chunks.append(_gfm(rows))
        chunks.append("")
    return "\n".join(chunks)


def extract(path: Path) -> dict:
    ok, msg = available()
    if not ok:
        return {"available": False, "body": "", "error": msg, "extra": {}}

    ext = path.suffix.lower().lstrip(".")
    raw = path.read_bytes()
    try:
        if ext == "csv":
            body = _extract_csv(raw)
            method = "csv-stdlib"
        elif ext == "xlsx":
            body = _extract_xlsx(raw)
            method = "openpyxl"
        elif ext == "pptx":
            body = _extract_pptx(raw)
            method = "python-pptx"
        else:
            return {"available": True, "body": "", "error": "not_applicable",
                    "extra": {}}
    except Exception as e:
        return {"available": True, "body": "",
                "error": f"{type(e).__name__}: {e}", "extra": {}}

    return {"available": True, "body": body, "error": None,
            "extra": {"extraction_method": method}}
