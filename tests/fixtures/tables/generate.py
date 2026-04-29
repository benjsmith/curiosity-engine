#!/usr/bin/env python3
"""Generate binary fixtures for the table-extraction harness.

Run once to produce:
  chem-buffer.pdf
  prose.pdf
  scientific-table.pdf
  merged-headers.xlsx
  pptx-3x3.pptx

The .csv fixture is committed verbatim (no generation).

Generated artefacts ARE checked in; the harness does not need this
script's deps installed at test time. Only re-run when assertions
or fixture content change.

Usage:  python3 generate.py
"""
import csv
from pathlib import Path

HERE = Path(__file__).resolve().parent


def _gen_chem_buffer_pdf():
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Table, TableStyle, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors

    out = HERE / "chem-buffer.pdf"
    doc = SimpleDocTemplate(str(out), pagesize=letter)
    styles = getSampleStyleSheet()
    story = [
        Paragraph("Buffer composition for in vitro experiments", styles["Title"]),
        Spacer(1, 12),
        Paragraph(
            "The drug solution was prepared in phosphate-buffered saline as "
            "described in Table 1.",
            styles["BodyText"],
        ),
        Spacer(1, 12),
    ]
    data = [
        ["Component", "Concentration", "Role"],
        ["Na2HPO4·7H2O", "5.36 mM", "buffer"],
        ["KH2PO4", "1.06 mM", "buffer"],
        ["5-fluorouracil", "50 µM", "drug"],
        ["H2O", "balance", "solvent"],
    ]
    tbl = Table(data, colWidths=[160, 100, 80])
    tbl.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.black),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("ALIGN", (0, 0), (-1, -1), "LEFT"),
    ]))
    story.append(tbl)
    doc.build(story)


def _gen_prose_pdf():
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet

    out = HERE / "prose.pdf"
    doc = SimpleDocTemplate(str(out), pagesize=letter)
    styles = getSampleStyleSheet()
    story = [
        Paragraph("A prose-only control document", styles["Title"]),
        Spacer(1, 12),
        Paragraph(
            "Lorem ipsum dolor sit amet, consectetur adipiscing elit. Sed do "
            "eiusmod tempor incididunt ut labore et dolore magna aliqua. Ut "
            "enim ad minim veniam, quis nostrud exercitation ullamco laboris.",
            styles["BodyText"],
        ),
        Spacer(1, 12),
        Paragraph(
            "This document contains no tables. The text-extraction baseline "
            "should round-trip cleanly through pypdf with no special handling.",
            styles["BodyText"],
        ),
    ]
    doc.build(story)


def _gen_scientific_pdf():
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Table, TableStyle, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors

    out = HERE / "scientific-table.pdf"
    doc = SimpleDocTemplate(str(out), pagesize=letter)
    styles = getSampleStyleSheet()
    story = [
        Paragraph("Cytoskeletal protein expression in salt-stressed cells", styles["Title"]),
        Spacer(1, 12),
        Paragraph(
            "Cells were exposed to elevated [Na⁺] and [Mg²⁺] "
            "(see refs [1], [2]). α-tubulin and β-actin abundance "
            "was assayed by western blot.",
            styles["BodyText"],
        ),
        Spacer(1, 8),
    ]
    data = [
        ["Protein", "Stress", "Fold change"],
        ["α-tubulin", "[Na⁺]", "1.4"],
        ["β-actin", "[Mg²⁺]", "0.9"],
    ]
    tbl = Table(data, colWidths=[120, 100, 80])
    tbl.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.black),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
    ]))
    story.append(tbl)
    doc.build(story)


def _gen_merged_xlsx():
    import openpyxl
    from openpyxl.styles import Font

    out = HERE / "merged-headers.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Outcomes"
    ws["A1"] = "Trial outcomes (n=120)"
    ws.merge_cells("A1:D1")
    ws["A1"].font = Font(bold=True)
    headers = ["Group", "Dose (mg/kg)", "Survival (%)", "p"]
    for i, h in enumerate(headers, 1):
        ws.cell(row=2, column=i, value=h).font = Font(bold=True)
    rows = [
        ["Control", 0, 65, "-"],
        ["Low", 5, 78, 0.04],
        ["High", 25, 91, 0.001],
    ]
    for ri, r in enumerate(rows, 3):
        for ci, v in enumerate(r, 1):
            ws.cell(row=ri, column=ci, value=v)
    wb.save(out)


def _gen_pptx():
    from pptx import Presentation
    from pptx.util import Inches

    out = HERE / "pptx-3x3.pptx"
    prs = Presentation()
    blank = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank)
    title = slide.shapes.title
    title.text = "Architecture comparison"

    rows = [
        ["Model", "Params", "Tokens"],
        ["BERT", "110M", "3.3B"],
        ["GPT-2", "1.5B", "40B"],
        ["LLaMA-2", "7B", "2T"],
    ]
    table_shape = slide.shapes.add_table(
        rows=len(rows), cols=3,
        left=Inches(1), top=Inches(2),
        width=Inches(8), height=Inches(3),
    ).table
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            table_shape.cell(ri, ci).text = val
    prs.save(out)


def _gen_csv():
    out = HERE / "gene-expression.csv"
    rows = [
        ["gene", "fold_change", "p_value", "context"],
        ["TP53", "2.3", "0.001", "colorectal"],
        ["BRCA1", "-1.8", "0.02", "breast"],
        ["MYC", "4.5", "0.0001", "lymphoma"],
        ["KRAS", "3.1", "0.003", "pancreatic"],
        ["EGFR", "2.7", "0.005", "lung"],
    ]
    with out.open("w", newline="") as fh:
        w = csv.writer(fh)
        for r in rows:
            w.writerow(r)


def main():
    _gen_chem_buffer_pdf()
    _gen_prose_pdf()
    _gen_scientific_pdf()
    _gen_merged_xlsx()
    _gen_pptx()
    _gen_csv()
    print("Generated:")
    for p in sorted(HERE.iterdir()):
        if p.suffix in (".pdf", ".xlsx", ".pptx", ".csv"):
            print(f"  {p.name}  ({p.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
